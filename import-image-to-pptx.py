from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import os
from glob import glob

img_types = ('*.png', '*.jpg')
grabbed_imgs = []
for files in img_types:
    grabbed_imgs.extend(glob(files))

grabbed_imgs.sort()

pptx_file = glob("*.pptx")[0]

prs = Presentation(pptx_file)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


for index, each_img in enumerate(grabbed_imgs):
    im = Image.open(each_img)
    img_width, img_height = im.size
    left = (prs.slide_width - img_width) / 3.5
    top = Inches(0.3)
    height = Inches(6.8)

    # blank_slide = prs.slide_layouts[6]
    # slide = prs.slides.add_slide(blank_slide)
    slide = prs.slides[index]
    slide.shapes.add_picture(each_img, left, top, height=height)

prs.save(pptx_file)
os.startfile(pptx_file)